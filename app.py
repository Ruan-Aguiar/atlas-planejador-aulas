import streamlit as st
from groq import Groq
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import io
import json

# ── Chave da API ─────────────────────────────────────────────────────────────
load_dotenv()
try:
    api_key = st.secrets.get("GEMINI_API_KEY", None)
except:
    api_key = os.getenv("GEMINI_API_KEY")

client = Groq(api_key=api_key) if api_key else None

# ── Configuração da página ────────────────────────────────────────────────────
st.set_page_config(
    page_title="ATLAS – Planejador de Aulas",
    page_icon="📚",
    layout="centered"
)

# ── Estilos ───────────────────────────────────────────────────────────────────
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Sora:wght@400;600;700&display=swap');
    html, body, [class*="css"] { font-family: 'Sora', sans-serif; }

    .titulo { font-size: 2.2rem; font-weight: 700; color: #1F60D8; margin-bottom: 0; }
    .subtitulo { font-size: 1rem; color: #666; margin-bottom: 2rem; }

    .boas-vindas {
        background-color: #f4f7ff;
        border-radius: 12px;
        padding: 1.5rem 2rem;
        margin-bottom: 1.5rem;
        border-left: 4px solid #1F60D8;
        font-size: 0.97rem;
        color: #333;
        line-height: 1.7;
    }

    .area-selecionada {
        background-color: #e8f0fe;
        border-radius: 8px;
        padding: 0.5rem 1rem;
        font-size: 0.9rem;
        color: #1F60D8;
        font-weight: 600;
        display: inline-block;
        margin-bottom: 1rem;
    }

    .historico-item {
        background-color: #f9f9f9;
        border-left: 3px solid #1F60D8;
        border-radius: 6px;
        padding: 0.6rem 1rem;
        margin-bottom: 0.5rem;
        font-size: 0.88rem;
        color: #444;
    }

    .stButton > button {
        background-color: #1F60D8;
        color: white;
        font-weight: 600;
        border-radius: 8px;
        padding: 0.6rem 2rem;
        border: none;
        width: 100%;
        font-size: 1rem;
    }
    .stButton > button:hover { background-color: #1549b0; }
</style>
""", unsafe_allow_html=True)

# ── Prompts por área ──────────────────────────────────────────────────────────
INSTRUCOES_POR_AREA = {
    "Linguagens": "Foque em atividades de leitura, interpretação, produção textual e oralidade. Inclua sugestões de gêneros textuais relevantes. Estimule criatividade e expressão.",
    "Ciências Humanas": "Foque na análise crítica de contextos históricos, geográficos e sociais. Sugira fontes primárias e atividades de debate. Estimule o pensamento crítico e cidadania.",
    "Ciências da Natureza": "Inclua uma atividade prática ou experimental. Relacione o conteúdo ao cotidiano. Destaque o método científico. Sugira experimentos simples quando possível.",
    "Matemática": "Inclua exercícios com níveis variados de dificuldade. Apresente ao menos uma situação-problema contextualizada. Indique fórmulas e propriedades relevantes.",
    "Outra": "Adapte o planejamento de forma interdisciplinar e flexível. Foque nos objetivos pedagógicos centrais do tema."
}

INSTRUCOES_POR_DISCIPLINA = {
    "Português": "Use textos literários e não-literários como ponto de partida. Trabalhe gramática de forma contextualizada.",
    "Matemática": "Priorize resolução de problemas e raciocínio lógico. Use situações do dia a dia.",
    "Física": "Relacione conceitos com fenômenos observáveis. Inclua fórmulas e unidades de medida.",
    "Química": "Conecte reações químicas com situações cotidianas. Inclua equações quando pertinente.",
    "Biologia": "Use exemplos da natureza local. Inclua atividade de observação ou análise quando possível.",
    "História": "Contextualize os fatos históricos com o presente. Use fontes primárias e imagens.",
    "Geografia": "Use mapas, dados e análises espaciais. Conecte com questões ambientais e sociais.",
    "Filosofia": "Estimule o questionamento e o debate. Use textos filosóficos acessíveis.",
    "Sociologia": "Conecte conceitos com a realidade social dos alunos. Proponha análises críticas.",
    "Inglês": "Integre as quatro habilidades: reading, writing, listening e speaking.",
    "Educação Física": "Equilibre teoria e prática. Inclua atividades cooperativas e reflexão sobre corpo e saúde.",
    "Arte": "Conecte produção artística com análise e apreciação. Use diferentes linguagens.",
}

# ── Geração de texto via Groq ─────────────────────────────────────────────────
def chamar_groq(system_prompt, user_prompt, max_tokens=2000):
    resposta = client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        max_tokens=max_tokens
    )
    return resposta.choices[0].message.content

# ── Gerador de PPTX ───────────────────────────────────────────────────────────
VERDE   = RGBColor(0x1F, 0x60, 0xD8)
BRANCO  = RGBColor(0xFF, 0xFF, 0xFF)
CINZA   = RGBColor(0xF4, 0xF7, 0xFF)
ESCURO  = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT  = RGBColor(0x16, 0x49, 0xB0)

def nova_pptx():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    return prs

def slide_em_branco(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg_solid(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_txt(slide, texto, x, y, w, h, size=16, color=None, bold=False,
            align=PP_ALIGN.LEFT, italic=False):
    if color is None: color = ESCURO
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return box

def gerar_pptx_do_plano(tema, serie, area, disciplina, planejamento_texto):
    """Gera um PPTX estruturado a partir do texto do planejamento."""
    prs = nova_pptx()

    # ── Slide 1: Capa ────────────────────────────────────────────────────────
    s = slide_em_branco(prs)
    bg_solid(s, VERDE)
    add_rect(s, 0, 3.8, 10, 1.825, ACCENT)
    add_txt(s, "📚", 0.5, 0.5, 1.5, 1.5, size=54, color=BRANCO, align=PP_ALIGN.CENTER)
    add_txt(s, "Planejamento de Aula", 2.3, 0.7, 7.2, 0.6, size=18, color=RGBColor(0xA8, 0xC8, 0xFF))
    add_txt(s, tema.upper(), 2.3, 1.2, 7.2, 1.2, size=32, color=BRANCO, bold=True)
    add_txt(s, f"{serie}  ·  {disciplina if disciplina else area}", 2.3, 2.5, 7.2, 0.5, size=15, color=RGBColor(0xA8, 0xC8, 0xFF))
    add_txt(s, "ATLAS – Planejador de Aulas com IA  ·  Gerado automaticamente",
            0.3, 4.05, 9.4, 0.5, size=12, color=BRANCO, align=PP_ALIGN.CENTER)

    # ── Pedir ao Groq para estruturar o plano em seções ─────────────────────
    prompt_estrutura = f"""
Com base neste planejamento de aula, extraia e organize as informações em exatamente 5 seções:
1. Objetivos
2. Conteúdos
3. Metodologia
4. Recursos e Avaliação
5. Encerramento e Reflexão

Para cada seção, forneça de 2 a 4 tópicos curtos (máximo 15 palavras cada).
Responda APENAS em JSON válido, sem markdown, neste formato:
{{
  "objetivos": ["tópico 1", "tópico 2"],
  "conteudos": ["tópico 1", "tópico 2"],
  "metodologia": ["tópico 1", "tópico 2"],
  "recursos_avaliacao": ["tópico 1", "tópico 2"],
  "encerramento": ["tópico 1", "tópico 2"]
}}

Planejamento:
{planejamento_texto[:2000]}
"""
    try:
        raw = chamar_groq(
            "Você é um assistente que extrai informações de planejamentos pedagógicos e retorna JSON válido.",
            prompt_estrutura,
            max_tokens=800
        )
        raw = raw.strip()
        if raw.startswith("```"):
            raw = raw.split("```")[1]
            if raw.startswith("json"):
                raw = raw[4:]
        secoes = json.loads(raw.strip())
    except:
        secoes = {
            "objetivos": ["Desenvolver habilidades relacionadas ao tema"],
            "conteudos": ["Conteúdo conforme planejamento gerado"],
            "metodologia": ["Aula expositiva dialogada com atividades práticas"],
            "recursos_avaliacao": ["Recursos didáticos variados", "Avaliação formativa"],
            "encerramento": ["Síntese dos conteúdos", "Reflexão final com os alunos"]
        }

    # ── Slides de conteúdo ───────────────────────────────────────────────────
    config_slides = [
        ("🎯  Objetivos", "objetivos", VERDE),
        ("📖  Conteúdos", "conteudos", RGBColor(0x05, 0x6B, 0x8B)),
        ("🧪  Metodologia", "metodologia", RGBColor(0x1B, 0x6B, 0x3A)),
        ("🛠️  Recursos & Avaliação", "recursos_avaliacao", RGBColor(0x7B, 0x2D, 0x8B)),
        ("💬  Encerramento", "encerramento", RGBColor(0x8B, 0x45, 0x00)),
    ]

    for titulo_slide, chave, cor in config_slides:
        s = slide_em_branco(prs)
        bg_solid(s, RGBColor(0xF4, 0xF7, 0xFF))
        add_rect(s, 0, 0, 10, 1.15, cor)
        add_txt(s, titulo_slide, 0.4, 0.22, 9.2, 0.75, size=26, color=BRANCO, bold=True)
        add_txt(s, f"{tema}  ·  {serie}", 0.4, 1.28, 9.2, 0.38, size=12,
                color=RGBColor(0x66, 0x66, 0x66))

        itens = secoes.get(chave, [])
        add_rect(s, 0.3, 1.75, 9.4, 0.06, cor)

        for i, item in enumerate(itens[:4]):
            cy = 1.95 + i * 0.85
            add_rect(s, 0.3, cy, 0.08, 0.65, cor)
            add_rect(s, 0.3, cy, 9.4, 0.65, RGBColor(0xFF, 0xFF, 0xFF))
            add_rect(s, 0.3, cy, 0.08, 0.65, cor)
            add_txt(s, item, 0.55, cy + 0.12, 9.0, 0.5, size=14, color=ESCURO)

    # ── Slide final ──────────────────────────────────────────────────────────
    s = slide_em_branco(prs)
    bg_solid(s, VERDE)
    add_txt(s, "✅", 4.5, 1.0, 1.0, 1.0, size=48, color=BRANCO, align=PP_ALIGN.CENTER)
    add_txt(s, "Planejamento Pronto!", 1.0, 2.1, 8.0, 0.9,
            size=32, color=BRANCO, bold=True, align=PP_ALIGN.CENTER)
    add_txt(s, f"Tema: {tema}  ·  {serie}",
            1.0, 3.1, 8.0, 0.5, size=15,
            color=RGBColor(0xA8, 0xC8, 0xFF), align=PP_ALIGN.CENTER)
    add_txt(s, "Gerado pelo ATLAS – Planejador de Aulas com Inteligência Artificial",
            1.0, 4.5, 8.0, 0.5, size=12,
            color=RGBColor(0xA8, 0xC8, 0xFF), align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

def repaginar_pptx(arquivo_bytes, nome_arquivo):
    """Lê PPTX ou PDF, extrai texto, repagina via Groq e gera novo PPTX."""

    # Extração de texto
    texto_extraido = ""
    if nome_arquivo.endswith(".pptx"):
        prs_orig = Presentation(io.BytesIO(arquivo_bytes))
        for slide in prs_orig.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        linha = para.text.strip()
                        if linha:
                            texto_extraido += linha + "\n"
            texto_extraido += "\n---\n"
    else:
        st.error("Por enquanto apenas arquivos .pptx são suportados para repaginação.")
        return None

    if not texto_extraido.strip():
        st.error("Não foi possível extrair texto do arquivo.")
        return None

    # Pedir ao Groq para reorganizar em slides
    prompt_repaginar = f"""
Você receberá o conteúdo de uma apresentação educacional.
Reorganize este conteúdo em slides claros, dinâmicos e fáceis de entender para alunos do ensino médio.
Mantenha TODA a essência e o conteúdo técnico do professor.

Retorne APENAS JSON válido, sem markdown, neste formato:
{{
  "titulo": "Título da apresentação",
  "slides": [
    {{"titulo": "Título do slide", "topicos": ["tópico 1", "tópico 2", "tópico 3"]}},
    ...
  ]
}}

Gere entre 6 e 10 slides. Cada slide com no máximo 4 tópicos curtos (até 20 palavras cada).

Conteúdo da apresentação:
{texto_extraido[:3000]}
"""
    try:
        raw = chamar_groq(
            "Você é um especialista em design instrucional que reorganiza conteúdo educacional em apresentações claras.",
            prompt_repaginar,
            max_tokens=1500
        )
        raw = raw.strip()
        if raw.startswith("```"):
            raw = raw.split("```")[1]
            if raw.startswith("json"):
                raw = raw[4:]
        dados = json.loads(raw.strip())
    except Exception as e:
        st.error(f"Erro ao processar com IA: {e}")
        return None

    # Gerar novo PPTX
    prs = nova_pptx()
    cores_slides = [
        RGBColor(0x1F, 0x60, 0xD8),
        RGBColor(0x05, 0x6B, 0x8B),
        RGBColor(0x1B, 0x6B, 0x3A),
        RGBColor(0x7B, 0x2D, 0x8B),
        RGBColor(0x8B, 0x45, 0x00),
        RGBColor(0x1F, 0x60, 0xD8),
        RGBColor(0x05, 0x6B, 0x8B),
        RGBColor(0x1B, 0x6B, 0x3A),
        RGBColor(0x7B, 0x2D, 0x8B),
        RGBColor(0x8B, 0x45, 0x00),
    ]

    # Capa
    s = slide_em_branco(prs)
    bg_solid(s, VERDE)
    add_rect(s, 0, 3.8, 10, 1.825, ACCENT)
    add_txt(s, "📋", 0.5, 0.5, 1.5, 1.5, size=54, color=BRANCO, align=PP_ALIGN.CENTER)
    add_txt(s, dados.get("titulo", "Apresentação"), 2.3, 1.0, 7.2, 1.5,
            size=30, color=BRANCO, bold=True)
    add_txt(s, "Versão repaginada pelo ATLAS – mantendo a essência do professor",
            0.3, 4.05, 9.4, 0.5, size=12, color=BRANCO, align=PP_ALIGN.CENTER)

    # Slides de conteúdo
    for i, slide_data in enumerate(dados.get("slides", [])[:10]):
        s = slide_em_branco(prs)
        cor = cores_slides[i % len(cores_slides)]
        bg_solid(s, RGBColor(0xF4, 0xF7, 0xFF))
        add_rect(s, 0, 0, 10, 1.15, cor)
        add_txt(s, slide_data.get("titulo", f"Slide {i+1}"),
                0.4, 0.22, 9.2, 0.75, size=24, color=BRANCO, bold=True)
        add_rect(s, 0.3, 1.75, 9.4, 0.06, cor)

        topicos = slide_data.get("topicos", [])
        for j, topico in enumerate(topicos[:4]):
            cy = 1.95 + j * 0.85
            add_rect(s, 0.3, cy, 9.4, 0.65, RGBColor(0xFF, 0xFF, 0xFF))
            add_rect(s, 0.3, cy, 0.08, 0.65, cor)
            add_txt(s, topico, 0.55, cy + 0.12, 9.0, 0.5, size=14, color=ESCURO)

    # Slide final
    s = slide_em_branco(prs)
    bg_solid(s, VERDE)
    add_txt(s, "✅", 4.5, 1.0, 1.0, 1.0, size=48, color=BRANCO, align=PP_ALIGN.CENTER)
    add_txt(s, "Apresentação Repaginada!", 1.0, 2.1, 8.0, 0.9,
            size=30, color=BRANCO, bold=True, align=PP_ALIGN.CENTER)
    add_txt(s, "Conteúdo preservado · Visual renovado · Pronto para a aula",
            1.0, 3.1, 8.0, 0.5, size=14,
            color=RGBColor(0xA8, 0xC8, 0xFF), align=PP_ALIGN.CENTER)
    add_txt(s, "Gerado pelo ATLAS – Planejador de Aulas com Inteligência Artificial",
            1.0, 4.5, 8.0, 0.5, size=12,
            color=RGBColor(0xA8, 0xC8, 0xFF), align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

# ── Estado da sessão ──────────────────────────────────────────────────────────
if "area_selecionada" not in st.session_state:
    st.session_state.area_selecionada = None
if "historico" not in st.session_state:
    st.session_state.historico = []
if "planejamento_atual" not in st.session_state:
    st.session_state.planejamento_atual = None
if "dados_atuais" not in st.session_state:
    st.session_state.dados_atuais = {}

# ── Cabeçalho ─────────────────────────────────────────────────────────────────
st.markdown('<p class="titulo">📚 ATLAS</p>', unsafe_allow_html=True)
st.markdown('<p class="subtitulo">Planejador de Aulas com Inteligência Artificial</p>', unsafe_allow_html=True)
st.divider()

# ══════════════════════════════════════════════════════════════════════════════
# ABA DE NAVEGAÇÃO
# ══════════════════════════════════════════════════════════════════════════════
aba_plano, aba_repaginar, aba_historico = st.tabs([
    "📝 Gerar Planejamento",
    "🎨 Repaginar Apresentação",
    "📋 Histórico"
])

# ══════════════════════════════════════════════════════════════════════════════
# ABA 1 — GERAR PLANEJAMENTO
# ══════════════════════════════════════════════════════════════════════════════
with aba_plano:

    # Tela de seleção de área
    if st.session_state.area_selecionada is None:
        st.markdown("""
        <div class="boas-vindas">
            👋 <strong>Bem-vindo ao ATLAS!</strong><br><br>
            Gere planejamentos de aula completos e alinhados à BNCC em segundos.<br><br>
            Selecione a <strong>área do conhecimento</strong> da sua disciplina para começar:
        </div>
        """, unsafe_allow_html=True)

        areas = {
            "🗣️ Linguagens": "Linguagens",
            "🌍 Ciências Humanas": "Ciências Humanas",
            "🔬 Ciências da Natureza": "Ciências da Natureza",
            "📐 Matemática": "Matemática",
            "📌 Outra": "Outra"
        }

        col1, col2 = st.columns(2)
        for i, (label, valor) in enumerate(areas.items()):
            col = col1 if i % 2 == 0 else col2
            with col:
                if st.button(label, key=f"area_{label}"):
                    st.session_state.area_selecionada = valor
                    st.session_state.planejamento_atual = None
                    st.rerun()

    # Formulário de planejamento
    else:
        area = st.session_state.area_selecionada
        st.markdown(f'<div class="area-selecionada">📂 Área: {area}</div>', unsafe_allow_html=True)

        if st.button("← Trocar área", key="trocar_area"):
            st.session_state.area_selecionada = None
            st.session_state.planejamento_atual = None
            st.rerun()

        st.markdown("<br>", unsafe_allow_html=True)

        col1, col2 = st.columns(2)
        with col1:
            tema = st.text_input("Tema da aula *", placeholder="Ex: Fotossíntese")
            serie = st.text_input("Série / Ano *", placeholder="Ex: 7º ano do Ensino Fundamental")
        with col2:
            duracao = st.number_input("Duração (minutos) *", min_value=10, max_value=300, value=50, step=5)

            # Disciplinas por área
            disciplinas_por_area = {
                "Linguagens": ["Português", "Inglês", "Educação Física", "Arte", "Outra"],
                "Ciências Humanas": ["História", "Geografia", "Filosofia", "Sociologia", "Outra"],
                "Ciências da Natureza": ["Biologia", "Física", "Química", "Outra"],
                "Matemática": ["Matemática", "Outra"],
                "Outra": ["Outra"]
            }
            opcoes = disciplinas_por_area.get(area, ["Outra"])
            disciplina = st.selectbox("Disciplina", opcoes)

        habilidades = st.text_input(
            "Habilidades BNCC (opcional)",
            placeholder="Ex: EF07CI07 – Caracterizar os principais ecossistemas brasileiros"
        )

        st.markdown("<br>", unsafe_allow_html=True)

        if st.button("Gerar Planejamento", key="gerar_plano"):
            erros = []
            if not tema.strip(): erros.append("O campo **Tema da aula** é obrigatório.")
            if not serie.strip(): erros.append("O campo **Série / Ano** é obrigatório.")
            if not api_key: erros.append("Chave da API não encontrada. Verifique o arquivo **.env**.")

            if erros:
                for e in erros:
                    st.warning(e)
            else:
                with st.spinner("Gerando planejamento, aguarde..."):
                    instrucoes_area = INSTRUCOES_POR_AREA.get(area, "")
                    instrucoes_disc = INSTRUCOES_POR_DISCIPLINA.get(disciplina, "")

                    prompt = f"""
Monte um planejamento de aula detalhado com base nas seguintes infor